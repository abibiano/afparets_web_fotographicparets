#!/usr/bin/env python3
"""Build the FOTOgraphic Parets 2026 'Bodegons' presentation.
Replica fidelment l'estil de la presentació original (alternança crema/fosc,
tipografies Georgia/Verdana, capçalera daurada, índex amb cercles, etc.)
Inclou notes del presentador amb timing per a cada slide.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image as PILImage

FOTOS = Path("/Users/abibiano/Projects/Web/afparets_web_fotographicparets/assets/ponencia/fotos")
OUT = Path("/Users/abibiano/Projects/Web/afparets_web_fotographicparets/assets/ponencia/ponencia_bodegons_2026.pptx")

# ── Paleta original ──────────────────────────────────────────
CREAM = RGBColor(0xF5, 0xF0, 0xEA)     # fons clar
CREAM_2 = RGBColor(0xF0, 0xE8, 0xDD)   # crema secundari (cards)
DARK = RGBColor(0x1A, 0x14, 0x10)      # fons fosc / text fosc
DARK_2 = RGBColor(0x3A, 0x2E, 0x24)    # marró mig (cards en fosc)
GOLD = RGBColor(0xC4, 0x94, 0x5A)      # daurat principal
GOLD_DARK = RGBColor(0x8B, 0x6E, 0x4E) # daurat fosc / subtitles
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)      # text secundari sobre crema

# Tipografies (com l'original)
TITLE_FONT = "Georgia"
BODY_FONT = "Verdana"

# Mida (idem original: 10" × 5.625" = 16:9)
SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)

# ── Helpers ──────────────────────────────────────────────────

def new_pres():
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def _set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _progress_bar(slide, current, total):
    """Barra superior de 2 colors: daurat = consumit, fosc-suau = restant.
    Permet al públic veure el progrés només mirant la barra."""
    bar_h = Inches(0.10)
    consumed_w = int(SLIDE_W * current / total)
    # part consumida (daurat)
    if consumed_w > 0:
        consumed = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, consumed_w, bar_h)
        consumed.line.fill.background()
        consumed.fill.solid()
        consumed.fill.fore_color.rgb = GOLD
        consumed.shadow.inherit = False
    # part restant (marró suau, més apagat)
    remaining_w = SLIDE_W - consumed_w
    if remaining_w > 0:
        remaining = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           consumed_w, 0, remaining_w, bar_h)
        remaining.line.fill.background()
        remaining.fill.solid()
        remaining.fill.fore_color.rgb = GOLD_DARK
        remaining.shadow.inherit = False


def _footer_logos(slide, on_dark=False):
    """Logos AFParets + FOTOgraphic a la part inferior, respectant l'aspect ratio."""
    afp = "logo_afparets_white.png" if on_dark else "logo_afparets.png"
    foto = "logo_fotographic.png"
    # AFParets a l'esquerra — ratio ~0.93 (81x87) · alçada 0.45"
    try:
        afp_img = PILImage.open(FOTOS / afp)
        ratio = afp_img.size[0] / afp_img.size[1]
        h = Inches(0.45)
        w = int(h * ratio)
        slide.shapes.add_picture(str(FOTOS / afp), Inches(0.25), Inches(5.15), w, h)
    except Exception:
        pass
    # FOTOgraphic a la dreta — ratio ~2.44 · alçada 0.40"
    try:
        foto_img = PILImage.open(FOTOS / foto)
        ratio = foto_img.size[0] / foto_img.size[1]
        h = Inches(0.40)
        w = int(h * ratio)
        x = SLIDE_W - w - Inches(0.25)
        slide.shapes.add_picture(str(FOTOS / foto), x, Inches(5.18), w, h)
    except Exception:
        pass


# Comptador global per al càlcul de progrés. S'inicialitza a 0 i
# s'incrementa cada cop que s'afegeix un slide. TOTAL_SLIDES s'estableix
# al final i les barres s'actualitzen en un segon pass.
SLIDE_REFS = []  # (slide_obj, position) per actualitzar després


def blank_cream(p):
    layout = p.slide_layouts[6]
    s = p.slides.add_slide(layout)
    _set_bg(s, CREAM)
    SLIDE_REFS.append((s, "cream"))
    _footer_logos(s, on_dark=False)
    return s


def blank_dark(p):
    layout = p.slide_layouts[6]
    s = p.slides.add_slide(layout)
    _set_bg(s, DARK)
    SLIDE_REFS.append((s, "dark"))
    _footer_logos(s, on_dark=True)
    return s


def finalize_progress_bars():
    """Afegeix la barra de progrés a tots els slides un cop coneixem el total."""
    total = len(SLIDE_REFS)
    for i, (slide, kind) in enumerate(SLIDE_REFS, 1):
        _progress_bar(slide, i, total)


def add_text(slide, text, left, top, width, height, *, size=14, color=DARK,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, font=BODY_FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = [text] if isinstance(text, str) else text
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        r = para.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb


def add_image_fit(slide, path, left, top, width, height, *, credit=None,
                  credit_color=None):
    """Centra la imatge dins l'àrea i posa el crèdit JUST SOTA de la foto real."""
    try:
        img = PILImage.open(path)
        iw, ih = img.size
    except Exception:
        iw, ih = 16, 9
    # Si volem que el crèdit càpiga sota la foto sense sortir de l'àrea,
    # reservem 0.22" per al crèdit
    credit_h = Inches(0.22) if credit else 0
    avail_h = height - credit_h
    ratio = min(width / iw, avail_h / ih) if iw and ih else 1
    new_w = int(iw * ratio)
    new_h = int(ih * ratio)
    cx = left + (width - new_w) // 2
    cy = top + (avail_h - new_h) // 2
    pic = slide.shapes.add_picture(str(path), cx, cy, new_w, new_h)
    if credit:
        # Crèdit just sota la foto, no flotant sobre l'àrea reservada
        credit_y = cy + new_h + Emu(20000)
        credit_x = cx
        credit_w = new_w
        add_text(slide, credit, credit_x, credit_y, credit_w, credit_h,
                 size=8, color=credit_color or GOLD_DARK,
                 align=PP_ALIGN.RIGHT, italic=True, font=BODY_FONT)
    return pic


def add_card(slide, left, top, width, height, *, fill=CREAM_2, border=GOLD,
             accent_left=True):
    """Card rectangular amb vora esquerra accent."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.shadow.inherit = False
    if accent_left and border:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                         Emu(60000), height)
        accent.line.fill.background()
        accent.fill.solid()
        accent.fill.fore_color.rgb = border
        accent.shadow.inherit = False


def add_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def add_title_top(slide, title, *, color=DARK, size=24):
    add_text(slide, title, Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.55),
             size=size, bold=True, color=color, font=TITLE_FONT)


# ── Constructors d'slides ────────────────────────────────────

def slide_title(p):
    s = blank_dark(p)
    # 2 retrats circulars centrats
    add_image_fit(s, FOTOS / "perfil_toni_circle.png",
                  Inches(3.55), Inches(1.0), Inches(1.4), Inches(1.4))
    add_image_fit(s, FOTOS / "perfil_alex_circle.png",
                  Inches(5.05), Inches(1.0), Inches(1.4), Inches(1.4))
    # Noms sota
    add_text(s, "Toni Barbany", Inches(3.05), Inches(2.45), Inches(2.4), Inches(0.3),
             size=11, bold=True, color=CREAM, align=PP_ALIGN.CENTER, font=BODY_FONT)
    add_text(s, "Alex Bibiano", Inches(4.55), Inches(2.45), Inches(2.4), Inches(0.3),
             size=11, bold=True, color=CREAM, align=PP_ALIGN.CENTER, font=BODY_FONT)
    # Títol gran
    add_text(s, "Bodegons", Inches(0.6), Inches(2.95), Inches(8.8), Inches(0.9),
             size=54, bold=True, color=CREAM, align=PP_ALIGN.CENTER, font=TITLE_FONT)
    add_text(s, "De l'ull a l'estudi", Inches(0.6), Inches(3.85), Inches(8.8), Inches(0.55),
             size=22, color=GOLD, align=PP_ALIGN.CENTER, italic=True, font=TITLE_FONT)
    # Footer
    add_text(s, "FOTOgraphic Parets 2026 · 6 de juny",
             Inches(0.6), Inches(4.65), Inches(8.8), Inches(0.3),
             size=11, color=GOLD_DARK, align=PP_ALIGN.CENTER, italic=True, font=BODY_FONT)
    add_notes(s,
        "[Bloc 1 · Presentació · 5 min total]\n\n"
        "ALEX (obertura):\n"
        "Bon dia a tothom, i benvinguts al FOTOgraphic Parets 2026. Sóc l'Alex Bibiano, "
        "soci de l'AFP des de fa uns quants anys. Els últims temps m'he dedicat de manera "
        "intensa al bodegó — és el gènere on em sento més còmode i on he tingut millors "
        "resultats darrerament.\n\n"
        "Avui estem aquí perquè el FOTOgraphic gira al voltant del bodegó, i he tingut la sort "
        "de preparar aquesta ponència conjuntament amb el Toni Barbany, un referent del gènere "
        "amb qui comparteixo l'afició pels bodegons, la llum i els objectes.\n\n"
        "[girar-se cap a Toni]\n"
        "Toni, et presentes?"
    )


def slide_index(p, items):
    """items: list of (num, title, brief)"""
    s = blank_cream(p)
    add_text(s, "Índex", Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.6),
             size=26, bold=True, color=DARK, font=TITLE_FONT)
    # cada element: cercle daurat + títol + brief
    y = Inches(1.2)
    row_h = Inches(0.6)
    gap = Inches(0.12)
    for num, title, brief in items:
        # Cercle daurat
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y,
                                   Inches(0.55), Inches(0.55))
        circ.line.fill.background()
        circ.fill.solid()
        circ.fill.fore_color.rgb = GOLD
        circ.shadow.inherit = False
        # Número dins el cercle
        n = s.shapes.add_textbox(Inches(0.6), y, Inches(0.55), Inches(0.55))
        ntf = n.text_frame
        ntf.margin_left = Emu(0); ntf.margin_right = Emu(0)
        ntf.margin_top = Emu(0); ntf.margin_bottom = Emu(0)
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = ntf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
        r = para.add_run(); r.text = num
        r.font.name = TITLE_FONT; r.font.size = Pt(15); r.font.bold = True
        r.font.color.rgb = WHITE
        # Títol + brief al costat
        add_text(s, title, Inches(1.35), y - Inches(0.05),
                 Inches(7.8), Inches(0.35),
                 size=15, bold=True, color=DARK, font=TITLE_FONT)
        add_text(s, brief, Inches(1.35), y + Inches(0.28),
                 Inches(7.8), Inches(0.3),
                 size=10, color=GOLD_DARK, font=BODY_FONT)
        y += row_h + gap
    add_notes(s,
        "[Índex · 30 segons]\n"
        "Alex: 'Avui parlarem de 6 blocs en 90 minuts. Comencem amb què és un bodegó, "
        "passem per composició i llum, i tanquem amb un teaser de les demos. "
        "Si teniu preguntes, interrompeu sense por.'"
    )


def slide_section(p, num, title, subtitle, time_note):
    s = blank_dark(p)
    add_text(s, num, Inches(0.6), Inches(1.0), Inches(8.8), Inches(1.7),
             size=110, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
             font=TITLE_FONT)
    # línia daurada de separació
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.25), Inches(3.0),
                              Inches(1.5), Emu(20000))
    sep.line.fill.background()
    sep.fill.solid(); sep.fill.fore_color.rgb = GOLD
    sep.shadow.inherit = False
    add_text(s, title, Inches(0.6), Inches(3.25), Inches(8.8), Inches(0.7),
             size=28, bold=True, color=CREAM, align=PP_ALIGN.CENTER, font=TITLE_FONT)
    add_text(s, subtitle, Inches(0.6), Inches(4.0), Inches(8.8), Inches(0.4),
             size=12, color=GOLD_DARK, align=PP_ALIGN.CENTER, italic=True, font=BODY_FONT)
    add_notes(s, time_note)


def slide_text(p, title, *, body=None, bullets=None, footnote=None,
                notes="", dark_bg=False, subtitle=None):
    s = blank_dark(p) if dark_bg else blank_cream(p)
    fg = CREAM if dark_bg else DARK
    body_color = CREAM_2 if dark_bg else GRAY
    add_title_top(s, title, color=fg, size=24)
    # decoració: línia daurada curta sota el títol
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.85),
                              Inches(1.2), Emu(30000))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
    bar.shadow.inherit = False
    if subtitle:
        add_text(s, subtitle, Inches(0.6), Inches(0.92), Inches(8.8), Inches(0.3),
                 size=11, color=GOLD_DARK, italic=True, font=BODY_FONT)
    top_y = Inches(1.30) if subtitle else Inches(1.15)
    if body:
        add_text(s, body, Inches(0.6), top_y, Inches(8.8), Inches(3.5),
                 size=14, color=body_color, font=BODY_FONT)
    if bullets:
        tb = s.shapes.add_textbox(Inches(0.6), top_y, Inches(8.8), Inches(3.6))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.space_after = Pt(8)
            r = para.add_run()
            r.text = "•  " + b
            r.font.name = BODY_FONT
            r.font.size = Pt(14)
            r.font.color.rgb = body_color
    if footnote:
        add_text(s, footnote, Inches(0.6), Inches(4.65), Inches(8.8), Inches(0.35),
                 size=11, color=GOLD_DARK, italic=True, font=BODY_FONT)
    add_notes(s, notes)


def slide_image(p, title, image, credit, *, caption=None, notes="",
                dark_bg=True):
    s = blank_dark(p) if dark_bg else blank_cream(p)
    fg = CREAM if dark_bg else DARK
    add_title_top(s, title, color=fg, size=22)
    add_image_fit(s, FOTOS / image, Inches(0.6), Inches(0.95),
                  Inches(8.8), Inches(3.7), credit=credit,
                  credit_color=GOLD if dark_bg else GOLD_DARK)
    if caption:
        add_text(s, caption, Inches(0.6), Inches(4.75), Inches(8.8), Inches(0.35),
                 size=11, color=GOLD if dark_bg else GOLD_DARK,
                 align=PP_ALIGN.CENTER, italic=True, font=BODY_FONT)
    add_notes(s, notes)


def slide_two_images(p, title, im1, cr1, im2, cr2, *, caption=None,
                     notes="", dark_bg=True):
    s = blank_dark(p) if dark_bg else blank_cream(p)
    fg = CREAM if dark_bg else DARK
    add_title_top(s, title, color=fg, size=22)
    add_image_fit(s, FOTOS / im1, Inches(0.5), Inches(0.95),
                  Inches(4.45), Inches(3.7), credit=cr1,
                  credit_color=GOLD if dark_bg else GOLD_DARK)
    add_image_fit(s, FOTOS / im2, Inches(5.05), Inches(0.95),
                  Inches(4.45), Inches(3.7), credit=cr2,
                  credit_color=GOLD if dark_bg else GOLD_DARK)
    if caption:
        add_text(s, caption, Inches(0.6), Inches(4.75), Inches(8.8), Inches(0.35),
                 size=11, color=GOLD if dark_bg else GOLD_DARK,
                 align=PP_ALIGN.CENTER, italic=True, font=BODY_FONT)
    add_notes(s, notes)


def slide_diagram(p, title, diagram, body, notes=""):
    s = blank_cream(p)
    add_title_top(s, title, color=DARK, size=22)
    add_image_fit(s, FOTOS / diagram, Inches(0.3), Inches(0.95),
                  Inches(5.0), Inches(3.8))
    add_text(s, body, Inches(5.5), Inches(1.05), Inches(4.2), Inches(3.7),
             size=12, color=GRAY, font=BODY_FONT)
    add_notes(s, notes)


def slide_3_principles(p, title, items, footnote=None, notes=""):
    """items: list of 6 (number, name, desc)"""
    s = blank_cream(p)
    add_title_top(s, title, color=DARK, size=24)
    # 2 columnes, 3 files
    col_w = Inches(4.3)
    row_h = Inches(1.15)
    gap_x = Inches(0.15)
    gap_y = Inches(0.08)
    start_x = Inches(0.45)
    start_y = Inches(1.1)
    for idx, (num, name, desc) in enumerate(items):
        col = idx % 2
        row = idx // 2
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * (row_h + gap_y)
        add_card(s, x, y, col_w, row_h)
        add_text(s, f"{num}. {name}", x + Inches(0.2), y + Inches(0.08),
                 col_w - Inches(0.3), Inches(0.3),
                 size=12, bold=True, color=DARK, font=TITLE_FONT)
        add_text(s, desc, x + Inches(0.2), y + Inches(0.43),
                 col_w - Inches(0.3), Inches(0.7),
                 size=9, color=GRAY, font=BODY_FONT)
    if footnote:
        add_text(s, footnote, Inches(0.6), Inches(4.7), Inches(8.8), Inches(0.3),
                 size=10, color=GOLD_DARK, italic=True, align=PP_ALIGN.CENTER,
                 font=BODY_FONT)
    add_notes(s, notes)


def slide_comparison_3(p, title, items, footnote=None, notes=""):
    """items: list of 3 (col_title, image, credit, bullets)"""
    s = blank_cream(p)
    add_title_top(s, title, color=DARK, size=22)
    col_w = Inches(3.1)
    gap = Inches(0.1)
    start_x = Inches(0.4)
    for i, (col_title, image, credit, bullets) in enumerate(items):
        x = start_x + i * (col_w + gap)
        # capçalera gold
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.0),
                                  col_w, Inches(0.4))
        hdr.line.fill.background(); hdr.fill.solid()
        hdr.fill.fore_color.rgb = GOLD
        hdr.shadow.inherit = False
        add_text(s, col_title, x, Inches(1.04), col_w, Inches(0.32),
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font=TITLE_FONT)
        # cos crema
        body = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.4),
                                   col_w, Inches(3.3))
        body.line.fill.background(); body.fill.solid()
        body.fill.fore_color.rgb = CREAM_2
        body.shadow.inherit = False
        if image:
            add_image_fit(s, FOTOS / image, x + Inches(0.1), Inches(1.5),
                          col_w - Inches(0.2), Inches(1.5))
            add_text(s, credit, x + Inches(0.1), Inches(3.05),
                     col_w - Inches(0.2), Inches(0.22),
                     size=7, color=GOLD_DARK, align=PP_ALIGN.CENTER,
                     italic=True, font=BODY_FONT)
        if bullets:
            tb = s.shapes.add_textbox(x + Inches(0.15), Inches(3.3),
                                       col_w - Inches(0.3), Inches(1.45))
            tf = tb.text_frame; tf.word_wrap = True
            for j, b in enumerate(bullets):
                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.space_after = Pt(1)
                r = para.add_run(); r.text = "• " + b
                r.font.name = BODY_FONT; r.font.size = Pt(8)
                r.font.color.rgb = GRAY
    if footnote:
        add_text(s, footnote, Inches(0.6), Inches(4.85), Inches(8.8), Inches(0.3),
                 size=11, color=GOLD, italic=True, bold=True,
                 align=PP_ALIGN.CENTER, font=BODY_FONT)
    add_notes(s, notes)


def slide_placeholder(p, title, note_text, notes=""):
    s = blank_cream(p)
    add_title_top(s, title, color=DARK, size=22)
    # caixa central amb borde daurat
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5),
                              Inches(7), Inches(2.5))
    box.line.color.rgb = GOLD; box.line.width = Pt(1.5)
    box.fill.solid(); box.fill.fore_color.rgb = CREAM_2
    box.shadow.inherit = False
    add_text(s, "⏳  Pendent — foto de Toni Barbany",
             Inches(1.5), Inches(1.8), Inches(7), Inches(0.5),
             size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
             font=TITLE_FONT)
    add_text(s, note_text, Inches(1.8), Inches(2.5), Inches(6.4), Inches(1.4),
             size=11, color=GRAY, align=PP_ALIGN.CENTER, font=BODY_FONT)
    add_notes(s, notes)


def slide_qa(p):
    s = blank_dark(p)
    add_text(s, "Preguntes?", Inches(0.6), Inches(1.6), Inches(8.8), Inches(1.2),
             size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font=TITLE_FONT)
    add_text(s, "Toni Barbany · Alex Bibiano",
             Inches(0.6), Inches(3.0), Inches(8.8), Inches(0.4),
             size=14, color=CREAM, align=PP_ALIGN.CENTER, font=BODY_FONT)
    add_text(s, "Gràcies", Inches(0.6), Inches(4.0), Inches(8.8), Inches(0.4),
             size=12, color=GOLD_DARK, align=PP_ALIGN.CENTER, italic=True,
             font=BODY_FONT)
    add_notes(s,
        "[Q&A · 5 min · Alex]\n\n"
        "ALEX:\n"
        "I ara el millor: preguntes. Cinc minuts. Qui és el primer?\n\n"
        "[Q&A obert — 5 minuts]\n\n"
        "[tancament]\n\n"
        "ALEX:\n"
        "Gràcies a tots. Espero que aquesta estona us hagi donat eines per entendre "
        "millor el que fareu avui. Ara a rodar."
    )


# ============================================================
# Construcció de la presentació
# ============================================================

p = new_pres()

# ── PORTADA ──
slide_title(p)

# ── ÍNDEX ──
slide_index(p, [
    ("01", "Presentació", "Qui som i estructura del dia"),
    ("02", "Què és un bodegó i tendències actuals", "Definició, gèneres, 4 estils"),
    ("03", "Composició: principis i errors", "5 principis + live view"),
    ("04", "La llum: concepte i efecte", "4 direccions + live view"),
    ("05", "La imatge com a relat", "Objectes, simbolisme, narrativa"),
    ("06", "Tancament i preguntes", "Consells finals i debat obert"),
])

# ════ BLOC 1 ════════════════════════════════════════════════
slide_section(p, "01", "Presentació", "Qui som, esquema del dia",
    "[Bloc 1 · 5 min]\n"
    "Alex obre la ponència. Estructura: presentació pròpia, presentació Toni, "
    "esquema del dia. Mantenir to natural — no llegir, conversar."
)

slide_two_images(p,
    "Toni Barbany — vidre i llum posterior",
    "toni_01.jpeg", "Toni Barbany",
    "toni_02.jpeg", "Toni Barbany",
    caption="Imatges gràfiques amb voluntat pictòrica",
    notes=(
        "[Toni es presenta · 1 min]\n\n"
        "TONI:\n"
        "Hola a tothom. Sóc el Toni Barbany. Porto anys competint en diversos gèneres "
        "però el bodegó és on he desenvolupat un estil més propi: vidre, llum posterior, "
        "la tensió entre la forma transparent i el fons fosc. Imatges gràfiques amb "
        "voluntat pictòrica.\n\n"
        "Avui estaré aquí parlant amb l'Alex i, una mica més tard, faré una demostració "
        "pràctica en directe.\n\n"
        "[Toni comenta breument les dues imatges — 30 segons. No cal entrar en detall "
        "tècnic, és una presentació visual del seu estil.]"
    )
)

slide_two_images(p,
    "Alex Bibiano — tenebrisme amb narrativa",
    "alex_01.jpeg", "Alex Bibiano",
    "alex_02.jpeg", "Alex Bibiano",
    caption="Objectes suspesos, tensió compositiva, simbolisme",
    notes=(
        "[Alex mostra la seva obra · 30 s]\n\n"
        "[Alex comenta breument les dues imatges — 30 segons. Remarcar la tensió "
        "compositiva i el simbolisme (objectes suspesos, amenaça, fragilitat).]"
    )
)

slide_text(p, "Estructura del dia",
    bullets=[
        "Ponència teòrica (90 min) — ara",
        "Presentació dels sets de bodegó",
        "Rotació: meitat practica · meitat veu demos (Toni + Alex)",
        "Canvi de grup i mateixes demos per a l'altra meitat",
        "Tot el que captureu es pot presentar al concurs",
    ],
    footnote="No és una classe magistral. Si teniu preguntes, feu-les en qualsevol moment.",
    notes=(
        "[Esquema del dia · 2 min · Alex]\n\n"
        "ALEX:\n"
        "Mireu, us explico com serà el dia perquè sapigueu on som en cada moment.\n\n"
        "Durant la propera hora i mitja parlarem de teoria: què és un bodegó, com es "
        "composa, com funciona la llum. És purament conceptual, però anirem introduint "
        "alguns exemples pràctics que podreu veure en pantalla amb la nostra càmera "
        "connectada per live view.\n\n"
        "Un cop acabada la ponència, presentarem els sets disponibles. El grup es "
        "dividirà en dues meitats: mentre una practica als sets de bodegó, l'altra "
        "seguirà les dues demostracions pràctiques en seqüència — primer la del Toni "
        "amb vidres retroil·luminats, de zero a imatge, i tot seguit la meva, un "
        "bodegó en clau baixa amb llum contínua des de la composició fins a l'edició "
        "a Photoshop. Després es fa la rotació i l'altra meitat veu les mateixes "
        "demostracions. Tot el que captureu durant la jornada es pot presentar al concurs.\n\n"
        "[pausa breu]\n"
        "Una cosa: aquesta ponència no és una classe magistral. Si teniu preguntes "
        "mentre parlem, feu-les. Si no esteu d'acord amb alguna cosa que diem, "
        "digueu-ho. Preferim el debat al monòleg."
    )
)

# ════ BLOC 2 ════════════════════════════════════════════════
slide_section(p, "02", "Què és un bodegó i tendències actuals", "Definició, gèneres, 4 estils",
    "[Bloc 2 · 15 min · Alex amb intervencions de Toni]"
)

slide_text(p, "Definició",
    body=(
        "Fotografia d'objectes inanimats, disposats artísticament, "
        "amb intenció estètica, narrativa o simbòlica.\n\n"
        "Aparentment simple — però la definició té un matís important que sovint genera preguntes."
    ),
    notes=(
        "[Definició · 1 min · Alex]\n\n"
        "ALEX:\n"
        "Comencem per la pregunta bàsica: què és un bodegó?\n\n"
        "La definició acadèmica és: fotografia d'objectes inanimats, disposats artísticament, "
        "amb intenció estètica, narrativa o simbòlica. Però aquesta definició té un matís "
        "important que genera moltes preguntes.\n\n"
        "[Pausa abans del clic — el slide següent és el gat]"
    )
)

slide_image(p, "...però hi ha gats premiats",
    "still_life_with_cat.jpg",
    'Marina Krasovska · Letònia · "Still life with cat" · Medalla Dorada FIAP · La Rioja',
    caption="El que defineix el gènere no és l'absència de vida — és la composició, la llum i la intenció.",
    notes=(
        "[2 min · Alex]\n\n"
        "ALEX:\n"
        "Mireu aquesta imatge. Guanyadora d'un concurs internacional de bodegons. Hi ha un gat. "
        "Un gat viu, en moviment, clarament l'element central de la composició. Tècnicament no "
        "és un \"objecte inanimat\" — i, malgrat això, ningú dubta que és un bodegó.\n\n"
        "Per què? Perquè el que defineix el gènere no és la presència estricta d'objectes inerts, "
        "sinó l'ENFOCAMENT: la composició intencionada, la llum treballada, la narrativa visual, "
        "l'atmosfera. Aquí el gat és el protagonista d'una escena barroca perfecta — peix, "
        "cistelles, una tetera de plata, vegetals — composada i il·luminada com ho hauria fet "
        "un mestre del segle XVII.\n\n"
        "La definició \"inanimat\" és un punt de partida útil, no una norma rígida. Animals, "
        "flors fresques, fins i tot elements en moviment poden aparèixer en un bodegó si "
        "serveixen la composició i la intenció artística és clara."
    )
)

slide_two_images(p, "Orígens pictòrics: el barroc espanyol",
    "Bodegón de caza, hortalizas y frutas.jpg",
    'Juan Sánchez Cotán · "Bodegón con caza, verdura y fruta" · c.1602',
    "Bodegón con cacharros.jpg",
    'Francisco de Zurbarán · "Bodegón con cacharros" · c.1650',
    caption="La fotografia de bodegó hereta del barroc: llum dramàtica, fons fosc, objectes elevats a art.",
    notes=(
        "[2 min · Alex]\n\n"
        "ALEX:\n"
        "Mireu aquestes dues imatges. Sánchez Cotán, 1602. Zurbarán, cap al 1650. Si les veieu "
        "i penseu \"semblen fotos\", és perquè la fotografia de bodegó ha après conscientment "
        "d'ells: la llum dramàtica que ve d'un sol costat, el fons fosc, els objectes quotidians "
        "elevats a art per la manera com estan il·luminats. Res d'això ho va inventar la "
        "fotografia — ho vam prendre prestat de la pintura barroca.\n\n"
        "Però no tots els bodegons són iguals, i aquí és on sovint es produeix la confusió."
    )
)

slide_comparison_3(p, "On s'acaba el bodegó? — Fine art · Food · Producte",
    [
        ("Bodegó fine art", "Watering_Can_and_Radish_PM-copy.jpg", "Harold Ross",
         ["Expressió artística personal",
          "Narrativa, simbolisme, emoció",
          "Il·luminació intencionada",
          "Objectes triats per valor estètic",
          "Galeria o concurs"]),
        ("Food photography", "ella_olsson_food.jpeg", "Ella Olsson",
         ["L'apetit com a objectiu",
          "Colors vius i saturats",
          "Clau alta, llum natural",
          "Ingredients frescos i perfectes",
          "Editorial i xarxes socials"]),
        ("Fotografia de producte", "noah_erickson_phone.jpeg", "Noah Erickson",
         ["Vendre un producte concret",
          "Fons neutre i net",
          "Il·luminació uniforme",
          "Objecte nou i perfecte",
          "Catàleg i e-commerce"]),
    ],
    footnote="Si la teva imatge podria ser un anunci, probablement no és fine art.",
    notes=(
        "[Comparativa fine art / food / producte · 3 min · Alex + intervenció Toni]\n\n"
        "ALEX:\n"
        "Hi ha tres gèneres que semblen el mateix però no ho són gens.\n\n"
        "El bodegó fine art és expressió artística personal. L'objectiu és l'emoció, "
        "la narrativa, el simbolisme. Qui el mira hauria de sentir alguna cosa, no "
        "tenir gana ni voler comprar res.\n\n"
        "La food photography vol despertar l'apetit. Colors vius, clau alta, tot fresc "
        "i perfecte. És editorialista, és per a Instagram o una revista gastronòmica. "
        "Molt vàlida, però no és el que farem avui.\n\n"
        "I la fotografia de producte vol vendre un objecte concret. Fons neutre, "
        "il·luminació uniforme, producte nou i immaculat. Si hi ha una marca visible "
        "i l'objectiu és el catàleg o l'e-commerce, estem en aquest territori.\n\n"
        "El test fàcil: si la teva imatge podria ser un anunci, probablement no és un "
        "bodegó fine art.\n\n"
        "[mirar a Toni]\n"
        "Toni, tu al circuit federatiu veus moltes imatges. Quina és l'errada més "
        "habitual en aquest punt?\n\n"
        "TONI: [resposta lliure sobre confusió de gèneres en concursos — 1 min aproximat]"
    )
)

slide_text(p, "Tendències actuals · 4 estils",
    bullets=[
        "Tenebrisme / Clarobscur",
        "Minimal / Escandinau",
        "Conceptual / Surrealista",
        "Retorn al pictòric",
    ],
    footnote="L'important no és l'estil que tries — és triar-lo amb intenció.",
    notes=(
        "[Intro tendències · 30 s · Alex]\n\n"
        "ALEX:\n"
        "Bé, suposant que volem fer bodegó fine art — que és el que teniu al concurs "
        "d'avui — quines tendències hi ha ara mateix?\n\n"
        "[clic per cadascun dels 4 estils següents]"
    )
)

slide_image(p, "Tenebrisme / Clarobscur",
    "Watering_Can_and_Radish_PM-copy.jpg",
    'Harold Ross · "Watering Can and Radish" · Light painting',
    caption="Fons fosc · objectes patinats · llum lateral dramàtica · hereu de Caravaggio i del barroc espanyol",
    notes=(
        "[Tenebrisme · 1.5 min · Alex]\n\n"
        "ALEX:\n"
        "El TENEBRISME — també conegut com a CLAROBSCUR o, en anglès, dark & moody — "
        "és el corrent dominant als concursos europeus ara mateix. Fons negres o molt "
        "foscos, objectes amb pàtina, llum lateral dura, molt contrast. Atmosfera "
        "d'estudi del segle XVII, directament hereva de Caravaggio i dels mestres "
        "espanyols que acabem de veure. Si busqueu \"bodegó fine art\" a Instagram, "
        "vuit de cada deu imatges són d'aquest estil."
    )
)

slide_image(p, "Minimal / Escandinau",
    "cerezas.jpg",
    'Leandro Cubillas Pérez · Espanya · "Cerezas" · Medalla Dorada FIAP 2026 · La Rioja',
    caption="Fons clar · un sol element · molt espai negatiu · llum suau",
    notes=(
        "[Minimal / Escandinau · 1 min · Alex]\n\n"
        "ALEX:\n"
        "El contrari: fons clars, un sol objecte, molt espai negatiu. Influència "
        "escandinava. Difícil de fer bé perquè no et pots amagar darrere de la "
        "complexitat — si l'objecte no és interessant o la llum no és perfecta, "
        "no hi ha res que ho amagui."
    )
)

slide_image(p, "Conceptual / Surrealista",
    "jaula-nube.jpeg",
    'Chema Madoz · "Sin título" · Premi Nacional de Fotografia 2000',
    caption="L'objecte ja és la idea — paradoxa, ironia, poesia visual",
    notes=(
        "[Conceptual / Surrealista · 1.5 min · Alex]\n\n"
        "ALEX:\n"
        "El CONCEPTUAL o SURREALISTA: bodegons on els objectes no s'arrangen, es "
        "col·loquen per crear una metàfora visual. La referència obligada és en "
        "Chema Madoz — Premi Nacional de Fotografia 2000. Objectes quotidians "
        "transformats per crear paradoxes, ironies, poesia. Mireu aquesta imatge: "
        "una gàbia normal i corrent, però en lloc d'un ocell hi ha un núvol atrapat. "
        "La idea és l'objecte. Sempre en blanc i negre, minimalista, sense "
        "post-producció digital. No busca bellesa decorativa — busca pregunta, "
        "sorpresa, lectura."
    )
)

slide_image(p, "Retorn al pictòric",
    "Paulette+Tavormina+Still+Life+with+Jamón+Ibérico+After+LM+2014.jpg",
    'Paulette Tavormina · "Still Life with Jamón Ibérico, after L.M." · 2014',
    caption='"After L.M." = a la manera de Luis Meléndez · barroc espanyol del XVIII en fotografia',
    notes=(
        "[Retorn al pictòric · 2 min · Alex + Toni + Alex]\n\n"
        "ALEX:\n"
        "I finalment, el que a mi m'interessa més: el retorn al pictòric. Buscar que "
        "la foto sembli una pintura. Aquest exemple és de Paulette Tavormina — el "
        "títol \"after L.M.\" vol dir \"a la manera de Luis Meléndez\", el mestre del "
        "bodegó espanyol del segle XVIII. Mireu l'analogia: si tornem al slide anterior "
        "de Sánchez Cotán i Zurbarán, veureu que aquesta fotografia podria haver estat "
        "pintada al segle XVII. Mateix fons negre, mateixos objectes (gerro de ceràmica, "
        "pa, embotit, alls), mateixa llum lateral suau. Però és una fotografia del 2014.\n\n"
        "[mirar a Toni]\n"
        "Toni, on et situes tu?\n\n"
        "TONI: [parla del seu estil — 45 segons]\n\n"
        "ALEX:\n"
        "Jo treballo sobretot el TENEBRISME amb voluntat narrativa — objectes amb "
        "pàtina, llum dramàtica, però amb una història al darrere. Però allò que "
        "m'importa és que l'estil sigui intencionat: que pugueu dir \"he triat aquest "
        "estil per aquesta raó\", no que us hagi quedat per casualitat.\n\n"
        "[pregunta al públic]\n"
        "A vosaltres, quin us atrau més? [pausa, mirar la sala]"
    )
)

# ════ BLOC 3 ════════════════════════════════════════════════
slide_section(p, "03", "Composició: principis i errors", "5 principis · live view · 25 min",
    "[Bloc 3 · 25 min · Alex protagonista, live view]"
)

slide_3_principles(p, "5 principis fonamentals", [
    ("1", "Protagonista clar",
     "Un sol centre de gravetat. La resta hi donen suport, no competeixen."),
    ("2", "Estructura geomètrica",
     "Triangle, diagonal o corba. Eviteu l'alineació horitzontal plana."),
    ("3", "Profunditat en Z",
     "Separa objectes en primer pla, pla mitjà i fons. Penseu Z."),
    ("4", "Solapament controlat",
     "Quan dos objectes se solapen, l'ull interpreta profunditat."),
    ("5", "Espai negatiu actiu",
     "40-60% buit no és farciment: és respiració, atmosfera."),
    ("✓", "Valida en B/N",
     "Si funciona sense color, la composició és sòlida."),
],
    footnote="Una composició sòlida aguanta fins i tot amb una llum mediocre.",
    notes=(
        "[Composició: overview · 3 min · Alex]\n\n"
        "ALEX:\n"
        "Ara entrem en la part que, al meu entendre, és la més important i la que més "
        "s'oblida: la composició.\n\n"
        "La gent s'obsessiona amb la llum, amb la tècnica, amb el postprocessat. I tot "
        "això importa. Però una imatge mal composada no la salva cap llum. En canvi, "
        "una composició sòlida aguanta fins i tot amb una llum mediocre.\n\n"
        "Cinc principis. Ràpids però fonamentals."
    )
)

slide_text(p, "1 · Protagonista clar",
    body=("Cada imatge necessita UN sol centre de gravetat. Un objecte principal "
          "al qual la vista torna sempre que entra al frame.\n\n"
          "La resta hi donen suport, no competeixen. Si dos objectes tenen el mateix "
          "pes visual, la imatge no té centre i l'ull marxa."),
    notes=(
        "[Principi 1 · 2 min · Alex]\n\n"
        "ALEX:\n"
        "U. Protagonista clar.\n\n"
        "Cada imatge necessita un sol centre de gravetat. Un objecte principal al qual "
        "la vista va sempre que entra al frame. La resta hi donen suport, no competeixen.\n\n"
        "Si teniu dos objectes amb el mateix pes visual, la imatge no té centre. L'ull "
        "no sap on posar-se i acaba marxant. Definiu el protagonista abans de posar res "
        "més al set."
    )
)

slide_text(p, "2 · Estructura geomètrica",
    body=("Organitzeu en triangle, diagonal o corba. Eviteu l'alineació horitzontal "
          "plana — tots els objectes en fila, com una foto de carnet.\n\n"
          "El triangle és el més versàtil: un objecte alt, un de mitjà, un de baix. "
          "Ja teniu estructura."),
    notes=(
        "[Principi 2 · 2 min · Alex]\n\n"
        "ALEX:\n"
        "Dos. Estructura geomètrica.\n\n"
        "Organitzeu els objectes en un triangle, una diagonal o una corba. Eviteu "
        "l'alineació horitzontal plana — tots els objectes en fila, a la mateixa "
        "alçada, com si fossin per a una foto de carnet. Avorrit i sense profunditat.\n\n"
        "El triangle és el més versàtil: un objecte alt, un de mitjà, un de baix. "
        "Ja teniu estructura."
    )
)

slide_text(p, "3 · Profunditat en Z",
    body=("Una foto és bidimensional però ha de semblar tridimensional.\n\n"
          "Separeu els objectes en plans: primer pla (Z1), pla mitjà (Z2), fons (Z3). "
          "Si tots estan a la mateixa distància de la càmera, la imatge s'aplana."),
    notes=(
        "[Principi 3 · 2 min · Alex]\n\n"
        "ALEX:\n"
        "Tres. Profunditat en Z.\n\n"
        "Una foto és bidimensional, però ha de semblar que té tres dimensions. Per "
        "aconseguir-ho necessiteu separar els objectes en plans: primer pla, pla mitjà "
        "i fons. Penseu en Z: davant, al mig, al darrere.\n\n"
        "Si tots els objectes estan a la mateixa distància de la càmera, l'aplanen. "
        "Separeu-los en profunditat."
    )
)

slide_text(p, "4 · Solapament controlat",
    body=("Quan dos objectes se solapen parcialment, la imatge guanya profunditat: "
          "el cervell interpreta 'aquest és davant d'aquell'.\n\n"
          "En canvi, una tangència — dos objectes que es toquen sense solapar-se — "
          "crea tensió visual i aplana la imatge."),
    footnote="→ A continuació, demostració en directe amb la càmera connectada (LIVE VIEW).",
    notes=(
        "[Principi 4 + LIVE VIEW · 3 min text + 5 min demo · Alex/Toni]\n\n"
        "ALEX:\n"
        "Quatre. Solapament controlat.\n\n"
        "Quan dos objectes se solapen parcialment — quan un tapa una mica de l'altre — "
        "la imatge guanya profunditat. El cervell interpreta \"aquest és davant d'aquell\" "
        "i crea la sensació de tres dimensions.\n\n"
        "En canvi, quan dos objectes es toquen però no se solapen — quan hi ha una "
        "tangència — es crea una tensió visual incòmoda i la imatge s'aplana. Ara us ho "
        "mostro en directe.\n\n"
        "→ LIVE VIEW · [Alex activa la càmera connectada, apareix el set en pantalla]\n\n"
        "[Alex col·loca dues ampolles que es toquen per la vora — tangència]\n"
        "Mireu: les ampolles es toquen però cap no tapa l'altra. Sembla que estan "
        "enganxades. Provoca una tensió incòmoda, no? Això és una tangència. La solució "
        "és senzilla.\n\n"
        "[Alex separa les ampolles 3 cm]\n"
        "Separades. Ara respiren.\n\n"
        "[Alex les apropa fins que una solapa clarament la part inferior de l'altra]\n"
        "O solapades. Ara una és davant de l'altra. La imatge guanya profunditat "
        "immediatament. Veieu la diferència?\n\n"
        "TONI:\n"
        "Quan reviso imatges per al jurat, les tangències es veuen a l'instant. I el "
        "problema és que quan estàs composant, amb el cap ficat al visor, costa veure-les. "
        "Per això és important allunyar-se de tant en tant i mirar el frame des de fora."
    )
)

slide_text(p, "5 · Espai negatiu actiu",
    body=("L'espai buit no és un problema a solucionar — és un recurs expressiu.\n\n"
          "Un 40-60% de l'enquadrament buit dóna respiració, aïlla el protagonista "
          "i crea atmosfera. El reflex instintiu és omplir-lo. Resistiu-ho."),
    notes=(
        "[Principi 5 · 2 min · Alex]\n\n"
        "ALEX:\n"
        "Cinc. Espai negatiu actiu.\n\n"
        "L'espai buit no és un problema a solucionar. És un recurs expressiu. Un 40-60% "
        "de l'enquadrament buit dóna respiració a la imatge, aïlla el protagonista i "
        "crea atmosfera.\n\n"
        "El reflex instintiu quan veiem una foto \"buida\" és omplir-la. Resistiu-ho. "
        "El buit és intencionat."
    )
)

slide_text(p, "Regla final · valida amb blanc i negre",
    body=("Canvieu la càmera a blanc i negre al live view (o feu-ho mentalment) "
          "i feu una foto de prova.\n\n"
          "Si la imatge funciona sense color — si la forma, la llum i la jerarquia "
          "estan clares — la composició és sòlida.\n\n"
          "Si no funciona en blanc i negre, no ho solucionarà el color."),
    notes=(
        "[Regla final + LIVE VIEW continuació · 2 min text + 4 min demo · Alex]\n\n"
        "ALEX:\n"
        "Una regla que uso sempre per validar una composició: canvieu la càmera a blanc "
        "i negre en el live view i feu una foto de prova. Si la imatge funciona sense "
        "color — si la forma, la llum i la jerarquia estan clares — la composició és "
        "sòlida. El color vindrà a sobre. Si no funciona en blanc i negre, no ho "
        "solucionarà el color.\n\n"
        "→ LIVE VIEW · errors finals:\n\n"
        "Dos errors més que us puc mostrar en directe.\n\n"
        "[Alex inclina lleugerament la càmera]\n"
        "Càmera no anivellada. Mireu la vora de la taula — queda inclinada. Sembla un "
        "error petit però es detecta immediatament i desvia l'atenció. A les càmeres "
        "teniu la quadrícula i el nivell electrònic. Feu-los servir.\n\n"
        "[Alex torna a anivellar]\n"
        "Millor.\n\n"
        "[Alex col·loca un objecte apuntant directament a la càmera]\n"
        "I la compressió per perspectiva: quan un objecte apunta directament cap a la "
        "càmera, la seva forma es comprimeix. Aquesta ampolla hauria de tenir una forma "
        "elegant — i no se li veu. La solució és simple.\n\n"
        "[Alex gira l'ampolla 30-45°]\n"
        "Trenta o quaranta-cinc graus. Ara es veu la seva silueta. L'objecte \"existeix\" "
        "dins la imatge."
    )
)

# ════ BLOC 4 ════════════════════════════════════════════════
slide_section(p, "04", "La llum: concepte i efecte", "4 direccions + live view · 23 min",
    "[Bloc 4 · 23 min · Toni inicia, Alex complementa + live view ràpid]"
)

slide_text(p, "La llum és narrativa, no tècnica",
    body=("Mateixos objectes. Mateix set. Mateixa càmera. Moveu la font de llum "
          "i teniu una imatge completament diferent.\n\n"
          "La llum decideix quin estat d'ànim transmet la foto, quina textura es veu, "
          "quins volums emergen i quins desapareixen."),
    footnote="Tècnica és com encendre la llum. Narrativa és on la poses.",
    notes=(
        "[Intro llum · 2 min · Toni]\n\n"
        "TONI:\n"
        "Ara parlem de llum. I vull que quedi clar des del principi: la llum no és una "
        "qüestió tècnica, és una qüestió narrativa.\n\n"
        "Teniu els mateixos objectes, el mateix set, la mateixa càmera. Moveu la font "
        "de llum i teniu una imatge completament diferent. La llum decideix quin estat "
        "d'ànim transmet la foto, quina textura es veu, quins volums emergen i quins "
        "desapareixen.\n\n"
        "Tècnica és com encendre la llum. Narrativa és on la poses."
    )
)

slide_text(p, "4 direccions bàsiques",
    bullets=[
        "Lateral dura (90°) — drama, textura, ombres tallants",
        "Lateral suau (softbox) — elegància, volum modelat",
        "Frontal — aplana, eviteu-la sense raó expressiva",
        "Contralum — separa del fons, ideal per a transparències",
    ],
    notes=(
        "[Overview 4 direccions · 1 min · Alex]\n\n"
        "ALEX:\n"
        "Quatre direccions bàsiques. [clic per cadascuna]"
    )
)

slide_diagram(p, "Llum lateral dura", "diagram_lateral_dura.png",
    "Font petita i sense difusor, col·locada al costat del subjecte.\n\n"
    "• Ombres llargues i tallants\n"
    "• Contrast alt\n"
    "• Revela cada imperfecció de la superfície\n\n"
    "Funciona molt bé amb fusta rugosa, metall oxidat, ceràmica amb gra.",
    notes=(
        "[Lateral dura · 2 min · Alex]\n\n"
        "ALEX:\n"
        "Llum lateral dura a 90 graus. La font queda al costat del subjecte, sense cap "
        "difusor. Ombres llargues i tallants, molt contrast, revela cada imperfecció de "
        "la superfície. Dramàtica. Funciona molt bé amb objectes de textura forta: fusta "
        "rugosa, metall oxidat, ceràmica amb gra."
    )
)

slide_diagram(p, "Llum lateral suau", "diagram_lateral_suau.png",
    "Mateixa posició però amb un difusor davant: softbox, tul, paper vegetal.\n\n"
    "• Ombres graduals\n"
    "• El volum es modela amb elegància\n"
    "• La textura segueix visible però no crida\n\n"
    "La més versàtil per al bodegó. Podríeu fer-ne tota la vida i no seria un error.",
    notes=(
        "[Lateral suau · 2 min · Alex]\n\n"
        "ALEX:\n"
        "La mateixa posició lateral però amb un difusor davant: softbox, tul, un paper "
        "vegetal. Les ombres es suavitzen, el volum es modela amb més elegància i la "
        "textura segueix visible però no crida. El més versàtil per al bodegó. Podríeu "
        "fer tota la vida bodegons amb aquesta llum i no seria un error."
    )
)

slide_diagram(p, "Llum frontal", "diagram_frontal.png",
    "La llum ve de la mateixa direcció que la càmera.\n\n"
    "• Aplana els volums\n"
    "• Les ombres cauen darrere, invisibles\n"
    "• Falta de relleu\n\n"
    "Eviteu-la, excepte si sabeu exactament què esteu fent i teniu una raó expressiva.",
    notes=(
        "[Frontal · 1.5 min · Alex]\n\n"
        "ALEX:\n"
        "Frontal. La llum ve de la mateixa direcció que la càmera. El problema: aplana. "
        "Les ombres cauen cap enrere, darrere dels objectes, invisibles. Tot queda pla "
        "i sense relleu. Eviteu-la, excepte si sabeu exactament el que esteu fent i "
        "teniu una raó expressiva."
    )
)

slide_diagram(p, "Contralum (back light)", "diagram_contralum.png",
    "La llum ve de darrere del subjecte, cap a la càmera.\n\n"
    "• Separa els objectes del fons\n"
    "• Crea una vora lluminosa\n"
    "• Genera profunditat i drama\n\n"
    "En vidres i objectes translúcids és especialment poderosa — la llum travessa el material.",
    notes=(
        "[Contralum · 2 min text + 1 min Toni · Alex/Toni]\n\n"
        "ALEX:\n"
        "Contralum: la llum ve de darrere del subjecte, cap a la càmera. Difícil de "
        "controlar però molt poderosa. Separa els objectes del fons, crea una vora "
        "lluminosa al voltant d'ells i genera profunditat. En vidres i objectes "
        "translúcids és especialment dramàtica perquè la llum travessa el material.\n\n"
        "[mirar a Toni]\n"
        "Toni, tu treballes molt amb contralum. Quin és el repte principal?\n\n"
        "TONI: [explica la seva experiència amb contralum — sense anticipar la demo "
        "amb vidres — 1 minut]"
    )
)

slide_text(p, "Live view · les 4 direccions i com afecten",
    bullets=[
        "Lateral dura (90°) — drama, ombres llargues, textura enriquida",
        "Lateral suau (difusor) — ombra suavitzada, volum modelat",
        "Frontal — aplanament total, sense relleu",
        "Contralum — silueta + vora lluminosa, separa del fons",
    ],
    footnote="→ LIVE VIEW: 1 objecte + 1 torxa. Veiem-ho directament a la pantalla.",
    notes=(
        "[Live view 4 direccions · 4-5 min · Alex (Toni comenta)]\n\n"
        "Setup necessari: càmera connectada a la pantalla (la mateixa del Bloc 3) + "
        "1 objecte simple (ampolla, gerro, fruita) + 1 torxa LED o panell petit.\n\n"
        "ALEX:\n"
        "Ara que hem vist els esquemes a la pissarra, ho fem en directe. Amb un sol "
        "objecte i una torxa, veureu com la mateixa escena canvia completament segons "
        "d'on ve la llum.\n\n"
        "→ LIVE VIEW\n\n"
        "[Alex agafa la torxa i la va movent al voltant del subjecte mentre la càmera "
        "projecta el resultat]\n\n"
        "1. LATERAL DURA (90°) → torxa al costat, sense difusor.\n"
        "   \"Mireu el drama: la textura s'enriqueix, però perdem informació al costat "
        "fosc.\"\n\n"
        "2. LATERAL SUAU → mateixa posició, però posem un paper vegetal entre la torxa "
        "i el subjecte.\n"
        "   \"Veieu la diferència? L'ombra es suavitza, el volum es modela amb "
        "elegància, però la textura es manté.\"\n\n"
        "3. FRONTAL → torxa al costat de la càmera.\n"
        "   \"Mireu què passa: tot s'aplana. No hi ha ombres a la vista, sembla una "
        "foto de carnet.\"\n\n"
        "4. CONTRALUM → torxa darrere del subjecte, apuntant cap a la càmera.\n"
        "   \"Mireu la silueta i la vora lluminosa. L'objecte se separa del fons sol.\"\n\n"
        "TONI:\n"
        "Aquí veieu en 30 segons el que els llibres expliquen en 10 pàgines. Triar la "
        "direcció de la llum és triar la història que voleu explicar."
    )
)

slide_diagram(p, "Qualitat de la llum", "diagram_qualitat.png",
    "Variable clau: la mida de la font en relació al subjecte.\n\n"
    "• Font petita i llunyana → llum DURA · ombres tallants · drama\n"
    "• Font gran i propera → llum DIFOSA · transicions subtils · elegància\n\n"
    "Podeu canviar la qualitat sense comprar res: poseu un paper vegetal entre la font i el subjecte.",
    notes=(
        "[Qualitat de la llum · 2 min · Alex]\n\n"
        "ALEX:\n"
        "A part de la direcció, hi ha la qualitat. I aquí la variable clau és la mida "
        "de la font de llum en relació al subjecte.\n\n"
        "Una font petita i llunyana dóna llum dura: ombres de vores tallants, molt "
        "contrast, molt drama.\n\n"
        "Una font gran i propera dóna llum suau: les ombres tenen gradació, les "
        "transicions de llum a ombra són subtils, l'ambient és més elegant.\n\n"
        "Podeu canviar la qualitat de la llum sense comprar res: acosteu la font al "
        "subjecte, o poseu una làmina de paper vegetal entre la font i el subjecte. "
        "Ja teniu difusor."
    )
)

slide_text(p, "El reflector",
    body=("Quan teniu una sola font de llum, el costat oposat del subjecte queda fosc. "
          "Poseu una cartolina blanca o platejada a l'altre costat i rebotarà llum cap "
          "a les ombres, suavitzant-les.\n\n"
          "No és una segona llum — és usar la mateixa llum dues vegades."),
    footnote="A la demo que faré després, el reflector tindrà un paper important.",
    notes=(
        "[El reflector · 2 min · Alex]\n\n"
        "ALEX:\n"
        "Una cosa ràpida que veig que molta gent no fa: el reflector.\n\n"
        "Quan teniu una sola font de llum, el costat oposat del subjecte queda fosc. "
        "Podeu posar una cartolina blanca o platejada a l'altre costat i rebotarà una "
        "mica de llum cap a les ombres, suavitzant-les. No és una segona llum — és "
        "usar la mateixa llum dues vegades.\n\n"
        "A la demo que faré després, el reflector tindrà un paper important. Ja ho veureu."
    )
)

slide_text(p, "Quines tècniques hi ha?",
    bullets=[
        "Flash d'estudi — potència alta, repetible, no veus el resultat fins després",
        "Llum contínua (LED, panells, tungstè) — WYSIWYG, ideal per aprendre",
        "Light painting — exposicions llargues, control absolut, requereix pràctica",
    ],
    footnote="Avui veureu la llum contínua aplicada a les dues demos — amb aproximacions ben diferents.",
    notes=(
        "[Tècniques · 1.5 min · Toni]\n\n"
        "TONI:\n"
        "Fins ara hem parlat de conceptes: direcció, qualitat, efecte. Però hi ha "
        "diverses maneres tècniques d'aplicar-ho.\n\n"
        "El flash d'estudi és la referència professional: potència alta, repetible, "
        "congela el moviment, molts modificadors disponibles. La limitació és que no "
        "veus el resultat fins que has disparat.\n\n"
        "La llum contínua — LED, panells, làmpades de tungstè — té l'avantatge que veus "
        "exactament el que obtens en temps real. El que veus és el que surti a la foto. "
        "Per aprendre, és imbatible.\n\n"
        "I el light painting és una qüestió a part: exposicions llargues a les fosques, "
        "il·lumines manualment cada zona amb una torxa o un panell LED. Control absolut, "
        "però requereix pràctica i temps.\n\n"
        "Avui veureu la llum contínua aplicada a les dues demos — amb dues aproximacions "
        "ben diferents."
    )
)

# ════ BLOC 5 ════════════════════════════════════════════════
slide_section(p, "05", "La imatge com a relat", "Objectes, simbolisme, narrativa · 15 min",
    "[Bloc 5 · 15 min · Toni protagonista, Alex modera]"
)

slide_text(p, "Selecció d'objectes",
    body=("No hi ha una llista de 'materials que funcionen' i 'que no funcionen' — això porta a imatges previstes.\n\n"
          "La pregunta correcta és: quin paper farà aquest objecte en la imatge?\n\n"
          "• Via clàssica: objectes amb temps damunt (metalls patinats, ceràmica artesanal, fusta envellida)\n"
          "• Via contemporània: objectes nets, fins i tot frescos, on el que crea l'interès és la composició i la sorpresa"),
    footnote="Si treus un objecte i la imatge no perd res, és que no hauria d'estar.",
    notes=(
        "[Selecció d'objectes · 4 min · Toni + Alex + Toni]\n\n"
        "TONI:\n"
        "Tot comença a l'objecte. Però aquí vull ser clar: no hi ha una llista de "
        "\"materials que funcionen\" i \"materials que no funcionen\". Això és una "
        "simplificació que porta a imatges previstes.\n\n"
        "El que realment cal preguntar-se és: quin paper farà aquest objecte en la "
        "imatge? Si la resposta és clara, qualsevol objecte pot funcionar.\n\n"
        "L'aproximació clàssica, la que jo tendeixo a fer, utilitza objectes que porten "
        "temps damunt: metalls patinats, ceràmica artesanal, fusta envellida. Aquests "
        "objectes \"parlen\" sols perquè ja porten una càrrega visual. La composició "
        "els acompanya.\n\n"
        "Però hi ha una altra via igual de vàlida: l'objecte contemporani, net, fins i "
        "tot fresc, on el que crea l'interès visual és la composició i la sorpresa. "
        "Tres cireres sobre un plat blanc. Una flor amb vapor. No hi ha pàtina, no hi "
        "ha història visible — però hi ha una decisió compositiva molt clara i un "
        "element inesperat que enganxa la mirada.\n\n"
        "[mirar a Alex]\n"
        "De fet, al darrer concurs internacional de bodegons de La Rioja, la foto "
        "guanyadora no era d'objectes antics. Era composició i sorpresa.\n\n"
        "ALEX:\n"
        "Exacte. El criteri no és l'edat de l'objecte — és la intenció. L'únic que cal "
        "evitar és l'objecte que no aporta res: ni per la seva forma, ni pel seu color, "
        "ni per la seva posició dins el frame. Si el treus i la imatge no perd res, "
        "és que no hauria d'estar.\n\n"
        "Toni, i quan tries objectes, com gestiones la paleta de color?\n\n"
        "TONI:\n"
        "La paleta ha de ser limitada. Dos o tres tons dominants, màxim. Si hi ha massa "
        "colors competint, la imatge es dispersa. Jo tendeixo als ocres, sienes i "
        "torrats — tons de terra, càlids. De tant en tant afegeixo un acent fred per "
        "crear tensió, però poc.\n\n"
        "Una cosa que ajuda: els colors càlids avancen òpticament cap a la càmera, els "
        "freds retrocedeixen. Podeu usar-ho per reforçar la profunditat sense moure cap "
        "objecte — col·loqueu el to càlid al primer pla i el fred al fons."
    )
)

slide_text(p, "Simbolisme i narrativa",
    body=("L'element que separa una bona foto d'un bodegó que t'agafa fort és la narrativa. "
          "Hi ha alguna cosa que passa entre els objectes?\n\n"
          "Tradició del *vanitas* barroc: crani, espelma apagada, rellotge, fruita marcida. "
          "Cada objecte és un símbol de la brevetat de la vida.\n\n"
          "Podeu fer el mateix sense cranis: una flor seca al costat d'una vela encesa. "
          "Un rellotge aturat. Un got buit."),
    footnote='"Quan la miro, em fa alguna cosa. Si la puc mirar 30 segons sense que em passi res, no hi ha relat."',
    notes=(
        "[Simbolisme i narrativa · 3 min · Toni + Alex + Toni]\n\n"
        "TONI:\n"
        "Però l'element que separa una bona foto d'un bodegó que t'agafa fort és la "
        "narrativa. Hi ha alguna cosa que passa entre els objectes?\n\n"
        "Penseu en la tradició del VANITAS barroc: el crani, l'espelma apagada, el "
        "rellotge, la fruita marcida. Cada objecte era un símbol. Un missatge sobre la "
        "brevetat de la vida. Podeu fer el mateix sense posar cranis: una flor seca al "
        "costat d'una vela encesa. Un rellotge aturat. Un got buit.\n\n"
        "La pregunta que em faig és: si trec aquest objecte, la foto perd alguna cosa? "
        "Si la resposta és no, el trec. Si la foto es veu incompleta sense ell, és que "
        "hi ha de ser.\n\n"
        "ALEX:\n"
        "I com saps quan una imatge \"té\" relat i quan no?\n\n"
        "TONI:\n"
        "Quan la miro i em fa alguna cosa. No necessàriament que m'agradi — que em faci "
        "alguna cosa. Que em provoqui una sensació, una pregunta, un record. Si la puc "
        "mirar trenta segons sense que em passi res, no hi ha relat."
    )
)

# Placeholders per a fotos de Toni
for i in range(1, 5):
    slide_placeholder(p, f"Una imatge meva ({i}/4)",
        "Imatge a substituir per la selecció final de Toni Barbany.\n"
        "Crèdit visible: \"Toni Barbany · [títol] · [any/premi]\"",
        notes=(
            f"[Foto Toni {i}/4 · {2 if i==1 else 1} min · Toni]\n\n"
            "TONI:\n"
            "Us mostro unes imatges meves i us explico les decisions darrere de cada una.\n\n"
            "[Toni explica cada imatge: per què aquells objectes, quina llum, per què "
            "aquell enquadrament — 4-5 minuts aproximats per al conjunt dels 4 slides]"
        )
    )

# ════ BLOC 6 ════════════════════════════════════════════════
slide_section(p, "06", "Tancament i preguntes", "Consells finals i debat obert",
    "[Bloc 6 · 10 min · Alex + Toni junts]"
)

slide_text(p, "4 consells pràctics per avui",
    bullets=[
        "Tria un sol protagonista abans de posar res més al set",
        "Composa en blanc i negre (al live view o mentalment)",
        "Mira el set des dels angles des dels quals NO faràs la foto",
        "Menys objectes, no més — si dubtes, treu-lo",
    ],
    notes=(
        "[4 consells · 2 min · Alex]\n\n"
        "ALEX:\n"
        "Estem a punt d'acabar. Quatre consells pràctics per quan estigueu rodant:\n\n"
        "U: Tria un sol protagonista abans de posar cap objecte més al set. Si no saps "
        "qui és el protagonista, no comencis a compondre.\n\n"
        "DOS: Recordeu la regla del blanc i negre que hem vist al Bloc de composició — "
        "apliqueu-la sempre abans de disparar. Si funciona sense color, la composició "
        "és sòlida.\n\n"
        "TRES: Mira el set des dels angles des dels quals NO faràs la foto. Rodeu el "
        "set, agafeu-lo en planta, mireu-lo des de baix. Us ajuda a entendre els volums "
        "i a detectar tangències que des del visor no veieu.\n\n"
        "QUATRE: Menys objectes, no més. Si dubteu entre treure o deixar un objecte, "
        "traieu-lo. Sempre podeu tornar-lo a posar. La tendència natural és acumular; "
        "el resultat habitualment és millor amb menys."
    )
)

slide_text(p, "El que vindrà tot seguit",
    bullets=[
        "Demo Toni — vidres retroil·luminats, perfils i siluetes amb cartolina negra",
        "Demo Alex — bodegó en clau baixa amb llum contínua + reflector + edició a Photoshop",
        "Sets oberts perquè rodeu el vostre bodegó",
        "Tot el que captureu es pot presentar al concurs FOTOgraphic 2026",
    ],
    footnote="Seran demostracions, no classes. Mireu les decisions; pregunteu si teniu dubtes.",
    notes=(
        "[Teaser de les demos · 2 min · Alex + Toni]\n\n"
        "ALEX:\n"
        "Per tancar, us explico el que vindrà a continuació.\n\n"
        "Primer, el Toni us farà una demostració pràctica de la seva tècnica de vidres "
        "retroil·luminats: ampolles i copes, llum posterior i cartolina negra darrere "
        "dels objectes per fer ressaltar els perfils i les siluetes. Una tècnica "
        "minimalista en equip però amb resultats molt impactants.\n\n"
        "Després, seré jo: muntaré un bodegó en clau baixa des de zero. Una sola font "
        "de llum contínua i un reflector. I quan tinguem la captura feta, farem una "
        "passada d'edició a Photoshop per tancar el procés complet.\n\n"
        "Seran demostracions, no classes. No heu d'apuntar paràmetres — simplement "
        "mireu les decisions, i pregunteu si teniu dubtes.\n\n"
        "[mirar a Toni]\n"
        "Alguna cosa que vulguis afegir, Toni?\n\n"
        "TONI: [afegeix el que consideri oportú — 30 segons]"
    )
)

slide_qa(p)


finalize_progress_bars()

p.save(OUT)
print(f"Generated: {OUT}")
print(f"Slides: {len(p.slides)}")
