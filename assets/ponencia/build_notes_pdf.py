#!/usr/bin/env python3
"""Notes PDF: each page = slide thumbnail (top) + presenter notes (bottom).
Per al Toni (no té PowerPoint) o per a impressió/revisió."""

import io
import re
from pathlib import Path
import fitz  # PyMuPDF
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from reportlab.lib.colors import HexColor
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph, Frame

PPTX = Path("/Users/abibiano/Projects/Web/afparets_web_fotographicparets/assets/ponencia/ponencia_bodegons_2026.pptx")
SLIDES_PDF = Path("/tmp/ponencia_preview.pdf")
OUT = Path("/Users/abibiano/Projects/Web/afparets_web_fotographicparets/assets/ponencia/ponencia_bodegons_2026_NOTES.pdf")

# Paleta original
BG = HexColor("#F5F0EA")
DARK = HexColor("#1A1410")
GOLD = HexColor("#C4945A")
GOLD_DARK = HexColor("#8B6E4E")
GRAY = HexColor("#555555")
WHITE = HexColor("#FFFFFF")

PAGE_W, PAGE_H = A4
MARGIN = 12 * mm

# Llegir notes
print("Llegint pptx...")
pres = Presentation(str(PPTX))
notes = [s.notes_slide.notes_text_frame.text.strip() for s in pres.slides]
print(f"  → {len(notes)} slides")

# Obrir PDF de slides
print("Obrint PDF de slides...")
slides_doc = fitz.open(str(SLIDES_PDF))
print(f"  → {slides_doc.page_count} pàgines")

# Generar PDF de notes
print(f"Generant {OUT}...")
c = canvas.Canvas(str(OUT), pagesize=A4)

notes_style = ParagraphStyle('Notes', fontName='Helvetica', fontSize=10,
                              textColor=DARK, leading=14, alignment=TA_LEFT)

for i, page in enumerate(slides_doc):
    slide_num = i + 1
    note_text = notes[i] if i < len(notes) else ""

    # Fons crema
    c.setFillColor(BG)
    c.rect(0, 0, PAGE_W, PAGE_H, fill=1, stroke=0)

    # Barra superior daurada
    c.setFillColor(GOLD)
    c.rect(0, PAGE_H - 7*mm, PAGE_W, 7*mm, fill=1, stroke=0)
    c.setFillColor(WHITE)
    c.setFont("Helvetica-Bold", 10)
    c.drawString(MARGIN, PAGE_H - 5.3*mm, f"Slide {slide_num} / {len(notes)}")
    c.drawRightString(PAGE_W - MARGIN, PAGE_H - 5.3*mm,
                       "Ponència Bodegons · FOTOgraphic Parets 2026")

    # Render diapositiva a 200 DPI
    pix = page.get_pixmap(matrix=fitz.Matrix(2.5, 2.5))
    img_data = pix.tobytes("png")

    # Àrea per a la diapositiva (a dalt, ~45% de l'alçada útil)
    avail_h = PAGE_H - 8*mm - MARGIN  # entre la barra superior i el final de pàgina
    slide_h = avail_h * 0.45
    slide_w = PAGE_W - 2 * MARGIN

    iw, ih = pix.width, pix.height
    ratio = min(slide_w / iw, slide_h / ih)
    new_w = iw * ratio
    new_h = ih * ratio
    cx = (PAGE_W - new_w) / 2
    cy = PAGE_H - 8*mm - 3*mm - new_h

    # Marc daurat fi
    c.setStrokeColor(GOLD_DARK)
    c.setLineWidth(0.4)
    c.rect(cx - 1, cy - 1, new_w + 2, new_h + 2, fill=0, stroke=1)
    c.drawImage(ImageReader(io.BytesIO(img_data)), cx, cy,
                width=new_w, height=new_h, mask='auto')

    # Separador i etiqueta
    sep_y = cy - 5*mm
    c.setStrokeColor(GOLD)
    c.setLineWidth(1.5)
    c.line(MARGIN, sep_y, PAGE_W - MARGIN, sep_y)
    c.setFont("Helvetica-Bold", 9)
    c.setFillColor(GOLD)
    c.drawString(MARGIN, sep_y - 4.5*mm, "NOTES DEL PRESENTADOR")

    # Notes (sota)
    notes_top = sep_y - 7*mm
    notes_bottom = MARGIN
    notes_h = notes_top - notes_bottom
    notes_w = PAGE_W - 2 * MARGIN

    if note_text.strip():
        # Parser: línies entre [...] o que comencen amb [ → cursiva daurada
        # ALEX:/TONI: → bold daurat
        # Línies amb [...] inline → la part entre [] en cursiva daurada
        html_parts = []
        for raw_line in note_text.split("\n"):
            line = raw_line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            stripped = line.strip()
            if not stripped:
                html_parts.append("&nbsp;")
            elif stripped.startswith("[") and stripped.endswith("]"):
                html_parts.append(f'<font color="#8B6E4E"><i>{line}</i></font>')
            elif stripped in ("ALEX:", "TONI:") or stripped.startswith(("ALEX:", "TONI:")):
                html_parts.append(f'<font color="#C4945A"><b>{line}</b></font>')
            else:
                # Resaltar [...] inline
                line = re.sub(r'(\[[^\]]+\])',
                              r'<font color="#8B6E4E"><i>\1</i></font>', line)
                html_parts.append(line)
        html = "<br/>".join(html_parts)
        para = Paragraph(html, notes_style)
        frame = Frame(MARGIN, notes_bottom, notes_w, notes_h,
                      leftPadding=4, bottomPadding=4, rightPadding=4, topPadding=4,
                      showBoundary=0)
        frame.addFromList([para], c)
    else:
        c.setFont("Helvetica-Oblique", 10)
        c.setFillColor(GOLD_DARK)
        c.drawCentredString(PAGE_W / 2, notes_top - 15*mm,
                             "(sense notes per a aquest slide)")

    c.showPage()
    if (i + 1) % 10 == 0:
        print(f"  → {i+1}/{len(notes)}")

c.save()
slides_doc.close()
print(f"✓ PDF generat: {OUT}")
print(f"  Mida: {OUT.stat().st_size // 1024} KB")
